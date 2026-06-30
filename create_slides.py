"""Generate CineMatch presentation slides (slides.pptx).

Run:
    python create_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x16, 0x21, 0x3e)
RED    = RGBColor(0xe6, 0x39, 0x46)
TEAL   = RGBColor(0x2a, 0x9d, 0x8f)
WHITE  = RGBColor(0xff, 0xff, 0xff)
YELLOW = RGBColor(0xf4, 0xa2, 0x61)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return tb


def accent_bar(slide):
    box(slide, Inches(0.5), Inches(1.15), Inches(12.33), Pt(3), fill_color=RED)


def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    txt(slide, title,
        Inches(0.5), Inches(0.18), Inches(12.33), Inches(0.9),
        size=32, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.5), Inches(0.9), Inches(12.33), Inches(0.3),
            size=13, color=RGBColor(0xbb, 0xbb, 0xcc), italic=True)


def bullets(slide, items, x, y, w, h, size=13, color=WHITE, marker="  >  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = marker + item
        run.font.size      = Pt(size)
        run.font.color.rgb = color


def kpi(slide, label, value, x, y, w=Inches(2.3), h=Inches(1.1)):
    box(slide, x, y, w, h, fill_color=NAVY,
        line_color=RGBColor(0x33, 0x44, 0x66), line_width=Pt(1))
    txt(slide, value, x, y + Inches(0.04), w, Inches(0.6),
        size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.62), w, Inches(0.42),
        size=11, color=RGBColor(0xaa, 0xaa, 0xcc), align=PP_ALIGN.CENTER)


def results_table(slide, x, y, col_w, row_h):
    headers = ["Model", "P@10", "NDCG@10", "HR@10", "Coverage", "Novelty", "ILD"]
    rows = [
        ("Random",       "0.003", "0.004", "0.032", "0.119", "14.27", "0.839"),
        ("MostPopular",  "0.064", "0.082", "0.337", "0.010",  "8.57", "0.952"),
        ("HighestRated", "0.021", "0.030", "0.164", "0.004", "10.76", "0.973"),
        ("ContentBased", "0.008", "0.009", "0.066", "0.321", "14.64", "0.201"),
        ("ItemItemCF",   "0.004", "0.004", "0.037", "0.059", "13.57", "0.815"),
        ("UserUserCF",   "0.010", "0.014", "0.088", "0.156", "12.59", "0.883"),
        ("MatrixFact",   "0.090", "0.112", "0.498", "0.068",  "9.55", "0.926"),
    ]
    best = ["0.090", "0.112", "0.498", "0.321", "14.64", "0.973"]

    for ci, h in enumerate(headers):
        bx = x + ci * col_w
        box(slide, bx, y, col_w, row_h, fill_color=NAVY)
        txt(slide, h, bx + Inches(0.05), y + Inches(0.05),
            col_w - Inches(0.1), row_h,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        row_y  = y + row_h * (ri + 1)
        is_mf  = row[0] == "MatrixFact"
        row_bg = (RGBColor(0x1e, 0x35, 0x55) if is_mf else
                  RGBColor(0x22, 0x28, 0x38) if ri % 2 == 0 else
                  RGBColor(0x1a, 0x1f, 0x30))
        for ci, val in enumerate(row):
            bx = x + ci * col_w
            box(slide, bx, row_y, col_w, row_h, fill_color=row_bg)
            is_best   = ci > 0 and val == best[ci - 1]
            val_color = TEAL if is_best else (WHITE if ci > 0 else RGBColor(0xdd, 0xdd, 0xff))
            txt(slide, val,
                bx + Inches(0.04), row_y + Inches(0.05),
                col_w - Inches(0.08), row_h,
                size=10, bold=is_best, color=val_color,
                align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

prs = new_prs()


# ── 1. TITLE ──────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x0d, 0x12, 0x26))
box(s, Inches(1.5), Inches(1.6), Inches(10.33), Inches(4.1), fill_color=NAVY)
txt(s, "CineMatch",
    Inches(1.7), Inches(1.8), Inches(9.93), Inches(1.2),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(4.5), Inches(3.0), Inches(4.33), Pt(4), fill_color=RED)
txt(s, "A Modular Movie Recommender System",
    Inches(1.7), Inches(3.15), Inches(9.93), Inches(0.5),
    size=20, color=RGBColor(0xcc, 0xcc, 0xee), align=PP_ALIGN.CENTER)
txt(s, "Recommender Systems  |  Individual Assignment  |  ESADE 2025/26",
    Inches(1.7), Inches(3.78), Inches(9.93), Inches(0.4),
    size=13, color=RGBColor(0x88, 0x88, 0xaa), align=PP_ALIGN.CENTER)
box(s, Inches(0), Inches(6.7), W, Inches(0.8), fill_color=RED)
txt(s, "7 algorithms   |   9 metrics at K=5,10,20   |   610 users   |   MovieLens Small",
    Inches(0), Inches(6.76), W, Inches(0.5),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── 2. PROJECT OVERVIEW ───────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Project Overview", "What we built and why")

pillars = [
    ("Dataset",
     "MovieLens Small\n100,836 ratings\n610 users  |  9,742 movies\n98.3% sparse\n22-year span (1996-2018)"),
    ("Algorithms",
     "4 families / 7 methods:\n\nBaselines (3)\nItem-Item CF\nUser-User CF\nContent-Based\nMatrix Factorization"),
    ("Evaluation",
     "9 metrics\nK = 5, 10 and 20\n\nTemporal holdout split\nRelevance >= 3.5 stars\n593 users evaluated"),
    ("Prototype UI",
     "Streamlit app\n3 interactive pages\n\nDataset Insights\nGet Recommendations\nCompare Models"),
]
col_w = Inches(3.0)
sx    = Inches(0.5)
for i, (title, body) in enumerate(pillars):
    bx = sx + i * (col_w + Inches(0.17))
    box(s, bx, Inches(1.45), col_w, Inches(5.6), fill_color=NAVY)
    box(s, bx, Inches(1.45), col_w, Inches(0.44), fill_color=RED)
    txt(s, title, bx + Inches(0.1), Inches(1.48), col_w - Inches(0.2), Inches(0.4),
        size=14, bold=True, color=WHITE)
    txt(s, body,  bx + Inches(0.12), Inches(2.0), col_w - Inches(0.24), Inches(4.8),
        size=12, color=RGBColor(0xcc, 0xcc, 0xff))


# ── 3. DATASET & EDA ─────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Dataset & EDA -- MovieLens Small",
            "9-section notebook (eda.ipynb) with executed outputs")

kpis_data = [("Users","610"), ("Movies","9,742"), ("Ratings","100,836"),
             ("Sparsity","98.3%"), ("Avg Rating","3.50")]
for i, (lbl, val) in enumerate(kpis_data):
    kpi(s, lbl, val, Inches(0.6) + i * Inches(2.48), Inches(1.45))

findings = [
    "Long tail: top 10% of movies receive ~50% of all ratings",
    "Ratings skew positive (mean 3.50) -- selection bias from voluntary watching",
    "~7% of users have < 20 ratings -- cold-start risk for CF and MF",
    "Tags cover ~45% of movies -- combined with genres in TF-IDF vocabulary",
    "22-year temporal span motivates a temporal (not random) train/test split",
    "Only 35.8% of items have >= 5 co-raters -- limits Item-Item CF recall",
]
txt(s, "Key Findings", Inches(0.5), Inches(2.82), Inches(12), Inches(0.35),
    size=15, bold=True, color=TEAL)
bullets(s, findings, Inches(0.5), Inches(3.22), Inches(12.3), Inches(3.8), size=14)


# ── 4. ALGORITHMS ────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "7 Algorithms -- 4 Families",
            "Uniform .fit() / .recommend() API across all methods")

algos = [
    ("Baselines (3)",        "baselines.py",
     "Random\nMostPopular\nHighestRated\n(min 20 ratings)",
     "No user history needed.\nLower/upper bounds for\npersonalised methods."),
    ("Item-Item CF",         "collaborative.py",
     "Cosine similarity\nbetween items\nmin_support=5\nk=20 neighbours",
     "score(u,i) =\nsum sim(i,j)*r(u,j)\n/ sum|sim(i,j)|"),
    ("User-User CF",         "collaborative.py",
     "Cosine similarity\nbetween users\nk=20 neighbours\nVectorised scoring",
     "score(u,i) =\nsum sim(u,n)*r(n,i)\n/ sum|sim(u,n)|"),
    ("Content-Based",        "content_based.py",
     "TF-IDF on genres\n+ user tags\nmax_features=5000\nCentred profile",
     "profile(u) =\nsum (r-mean_r)\n* vector(i)"),
    ("Matrix Fact.",         "matrix_factorization.py",
     "TruncatedSVD\n50 latent factors\nR ~ U*Sigma*Vt\nStored R-hat",
     "Captures latent\ntaste patterns\nacross all users"),
]
col_w = Inches(2.4)
sx    = Inches(0.3)
for i, (name, fname, details, formula) in enumerate(algos):
    bx = sx + i * (col_w + Inches(0.12))
    box(s, bx, Inches(1.45), col_w, Inches(5.6), fill_color=NAVY)
    box(s, bx, Inches(1.45), col_w, Inches(0.4), fill_color=RED)
    txt(s, name,    bx + Inches(0.08), Inches(1.47), col_w - Inches(0.16), Inches(0.38),
        size=12, bold=True, color=WHITE)
    txt(s, fname,   bx + Inches(0.08), Inches(1.93), col_w - Inches(0.16), Inches(0.28),
        size=9, color=YELLOW, italic=True)
    txt(s, details, bx + Inches(0.08), Inches(2.28), col_w - Inches(0.16), Inches(2.2),
        size=10, color=RGBColor(0xcc, 0xcc, 0xff))
    box(s, bx + Inches(0.08), Inches(4.52), col_w - Inches(0.16), Pt(1),
        fill_color=RGBColor(0x33, 0x44, 0x66))
    txt(s, formula, bx + Inches(0.08), Inches(4.62), col_w - Inches(0.16), Inches(2.1),
        size=9, color=TEAL, italic=True)


# ── 5. EVALUATION METHODOLOGY ────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Evaluation Methodology",
            "Temporal split + relevance threshold + 9 metrics at K=5,10,20")

box(s, Inches(0.5), Inches(1.45), Inches(5.8), Inches(5.6), fill_color=NAVY)
txt(s, "Train / Test Split", Inches(0.62), Inches(1.55), Inches(5.6), Inches(0.38),
    size=14, bold=True, color=TEAL)
split_body = (
    "Per-user temporal holdout:\n\n"
    "Ratings sorted by timestamp.\n"
    "Most recent 20% -> test set.\n"
    "Earlier 80% -> training set.\n\n"
    "Eliminates temporal leakage:\n"
    "a random shuffle would let future\n"
    "ratings 'inform' past training.\n\n"
    "Relevance threshold: rating >= 3.5\n"
    "A 0.5-star item no longer counts\n"
    "the same as a 5-star item.\n\n"
    "Train: 80,896   |   Test: 19,940\n"
    "Users evaluated: 593"
)
txt(s, split_body, Inches(0.62), Inches(2.0), Inches(5.55), Inches(4.8),
    size=12, color=WHITE)

metrics_groups = [
    ("Accuracy Metrics",
     ["Precision@K  -- fraction of recs that are relevant",
      "Recall@K  -- fraction of liked items retrieved",
      "NDCG@K  -- ranking quality, penalises late hits",
      "MRR@K  -- reciprocal rank of first hit",
      "Hit Rate@K  -- probability of >= 1 relevant rec"]),
    ("Beyond-Accuracy",
     ["Coverage  -- fraction of catalogue recommended",
      "Novelty  -- self-information (Zhou et al. 2010)",
      "Avg Popularity  -- popularity bias indicator",
      "ILD  -- intra-list diversity (Ziegler et al. 2005)"]),
]
rx, ry = Inches(6.6), Inches(1.45)
for group, items in metrics_groups:
    box(s, rx, ry, Inches(6.2), Inches(0.38), fill_color=RED)
    txt(s, group, rx + Inches(0.1), ry + Inches(0.04),
        Inches(6.0), Inches(0.32), size=13, bold=True, color=WHITE)
    ry += Inches(0.38)
    for item in items:
        box(s, rx, ry, Inches(6.2), Inches(0.47), fill_color=NAVY)
        txt(s, "  " + item, rx + Inches(0.08), ry + Inches(0.07),
            Inches(6.0), Inches(0.38), size=11, color=RGBColor(0xcc, 0xcc, 0xff))
        ry += Inches(0.47)
    ry += Inches(0.12)
txt(s, "Evaluated at K = 5, 10 and 20  --  results in results_multi_k.csv",
    Inches(6.6), ry + Inches(0.06), Inches(6.2), Inches(0.4),
    size=11, color=TEAL, italic=True)


# ── 6. RESULTS TABLE ─────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Results -- K=10, 593 Users, Relevance >= 3.5 Stars",
            "Temporal split  |  Teal = best per column")

results_table(s, Inches(0.4), Inches(1.35), Inches(1.7), Inches(0.56))
txt(s, "MatrixFact dominates all accuracy metrics.  ContentBased leads on novelty and coverage.  ILD reveals genre diversity per model.",
    Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.5),
    size=11, color=TEAL, italic=True)


# ── 7. MULTI-K COMPARISON ────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Rankings Hold Across K = 5, 10, 20",
            "P@K for all models -- results are not a K=10 artefact")

mk_rows = [
    ("Random",       "0.003", "0.003", "0.002", "Flat -- random noise"),
    ("MostPopular",  "0.079", "0.064", "0.054", "Decreasing -- popular items front-loaded"),
    ("HighestRated", "0.024", "0.021", "0.025", "Stable -- quality signal holds"),
    ("ContentBased", "0.010", "0.008", "0.006", "Decreasing -- best matches front-loaded"),
    ("ItemItemCF",   "0.003", "0.004", "0.004", "Stable (low throughout)"),
    ("UserUserCF",   "0.008", "0.010", "0.011", "Increasing -- signal spreads"),
    ("MatrixFact",   "0.103", "0.090", "0.079", "Decreasing -- dominates at all K"),
]
headers = ["Model", "P@5", "P@10", "P@20", "Trend"]
cw = Inches(1.6)
ct = Inches(4.3)
sx, sy = Inches(0.4), Inches(1.38)
rh = Inches(0.6)

for ci, h in enumerate(headers):
    bx = sx + ci * cw if ci < 4 else sx + 4 * cw
    w  = cw if ci < 4 else ct
    box(s, bx, sy, w, rh, fill_color=NAVY)
    txt(s, h, bx + Inches(0.05), sy + Inches(0.1), w - Inches(0.1), rh,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(mk_rows):
    ry    = sy + rh * (ri + 1)
    is_mf = row[0] == "MatrixFact"
    rbg   = (RGBColor(0x1e, 0x35, 0x55) if is_mf else
             RGBColor(0x22, 0x28, 0x38) if ri % 2 == 0 else
             RGBColor(0x1a, 0x1f, 0x30))
    for ci, val in enumerate(row):
        bx = sx + ci * cw if ci < 4 else sx + 4 * cw
        w  = cw if ci < 4 else ct
        box(s, bx, ry, w, rh, fill_color=rbg)
        vc = (TEAL if is_mf and 0 < ci < 4 else
              RGBColor(0xaa, 0xcc, 0xff) if ci == 4 else
              WHITE if ci > 0 else RGBColor(0xdd, 0xdd, 0xff))
        txt(s, val, bx + Inches(0.05), ry + Inches(0.1), w - Inches(0.1), rh,
            size=11, bold=(is_mf and 0 < ci < 4), color=vc,
            align=PP_ALIGN.CENTER if ci < 4 else PP_ALIGN.LEFT,
            italic=(ci == 4))


# ── 8. BEYOND-ACCURACY ───────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Beyond-Accuracy Analysis",
            "Coverage, Novelty, Popularity Bias, Intra-List Diversity (ILD)")

panels = [
    ("Catalog Coverage",
     "ContentBased: 32.1%  [best]\nUserUserCF:   15.6%\nRandom:       11.9%\nMatrixFact:    6.8%\nMostPopular:   1.0%  [worst]",
     "ContentBased covers 32x more\nof the catalogue. Each user's\nprofile points to a different\ncorner of the item space."),
    ("Novelty (self-info)",
     "ContentBased: 14.64  [best]\nRandom:        14.27\nItemItemCF:   13.57\nUserUserCF:   12.59\nMatrixFact:    9.55\nMostPopular:   8.57  [worst]",
     "MostPopular recommends\nwhat everyone already knows.\nContentBased surfaces\ngenuinely niche items."),
    ("Intra-List Diversity",
     "HighestRated: 0.973\nMostPopular:  0.952  [best div]\nMatrixFact:   0.926\nUserUserCF:   0.883\nItemItemCF:   0.815\nContentBased: 0.201  [lowest]",
     "ContentBased clusters recs\ntightly by genre (ILD=0.20).\nMMR re-ranking would help.\nPopular films span genres."),
    ("Popularity Bias (AvgPop)",
     "MostPopular: 216.7  [worst]\nMatrixFact:  126.3\nHighestRated: 66.1\nUserUserCF:   29.3\nRandom:       10.8\nContentBased:  9.6  [best]",
     "High accuracy models surface\npopular items -- less value\nfor users who know the\nclassics already."),
]
col_w = Inches(3.1)
sx    = Inches(0.3)
for i, (title, data, insight) in enumerate(panels):
    bx = sx + i * (col_w + Inches(0.12))
    box(s, bx, Inches(1.45), col_w, Inches(5.6), fill_color=NAVY)
    box(s, bx, Inches(1.45), col_w, Inches(0.38), fill_color=RED)
    txt(s, title, bx + Inches(0.08), Inches(1.47), col_w - Inches(0.16), Inches(0.36),
        size=12, bold=True, color=WHITE)
    txt(s, data,  bx + Inches(0.1), Inches(1.9), col_w - Inches(0.2), Inches(2.75),
        size=10, color=RGBColor(0xcc, 0xdd, 0xff))
    box(s, bx + Inches(0.08), Inches(4.68), col_w - Inches(0.16), Pt(1), fill_color=TEAL)
    txt(s, insight, bx + Inches(0.1), Inches(4.78), col_w - Inches(0.2), Inches(2.0),
        size=10, color=RGBColor(0x99, 0xff, 0xcc), italic=True)


# ── 9. ItemItemCF DEEP DIVE ──────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Item-Item CF: Why P@10 = 0.004",
            "A dataset sparsity constraint -- not a bug")

box(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(5.6), fill_color=NAVY)
txt(s, "The Sparsity Problem", Inches(0.62), Inches(1.55), Inches(5.8), Inches(0.38),
    size=14, bold=True, color=TEAL)
left_body = (
    "With 98.3% sparsity, most item\n"
    "pairs share very few co-raters.\n\n"
    "The min_support=5 filter:\n"
    "Items with < 5 raters have their\n"
    "similarities zeroed out. Without it,\n"
    "a single shared rater gives\n"
    "cosine_sim = 1.0 -- an artefact\n"
    "that floods the top-10 with scores\n"
    "equal to the user's max rating.\n\n"
    "With the filter:\n"
    "Only 35.8% of items survive.\n"
    "The other 64.2% score zero\n"
    "and cannot be ranked.\n\n"
    "min_support=2 was tested:\n"
    "P@10 dropped to 0.001 -- worse.\n"
    "Sparser items = noisier similarities."
)
txt(s, left_body, Inches(0.62), Inches(2.0), Inches(5.7), Inches(4.8),
    size=12, color=WHITE)

box(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.6), fill_color=NAVY)
txt(s, "What It Gets Right", Inches(6.92), Inches(1.55), Inches(5.8), Inches(0.38),
    size=14, bold=True, color=TEAL)
right_body = (
    "Where it does form opinions,\n"
    "the model behaves sensibly:\n\n"
    "AvgPopularity = 7.3\n"
    "-> Lowest of all personalised\n"
    "   methods. Surfaces genuinely\n"
    "   niche, long-tail items.\n\n"
    "Novelty = 13.57\n"
    "-> Above average. Recommendations\n"
    "   are non-obvious.\n\n"
    "ILD = 0.815\n"
    "-> Varied lists when it can\n"
    "   form an opinion.\n\n"
    "In production: pair with a\n"
    "popularity fallback for items\n"
    "below the support threshold.\n"
    "Standard hybrid design."
)
txt(s, right_body, Inches(6.92), Inches(2.0), Inches(5.7), Inches(4.8),
    size=12, color=WHITE)


# ── 10. STREAMLIT UI ─────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Interactive Prototype -- Streamlit",
            "3 pages, @st.cache_resource training, runs in the browser")

pages = [
    ("Page 1\nDataset Insights",
     ["5 KPI cards: Users, Movies,\n  Ratings, Sparsity, Avg Rating",
      "Rating distribution histogram",
      "Ratings-per-user distribution",
      "Top-15 most-rated movies",
      "Genre frequency bar chart",
      "Long-tail popularity curve (log)"]),
    ("Page 2\nGet Recommendations",
     ["Select any of 610 users",
      "Choose algorithm + list length",
      "User profile: top-rated movies\n  + favourite genres",
      "Cold-start warning if < 20 ratings",
      "MatrixFact shows predicted\n  score per recommendation",
      "Genre + popularity per item"]),
    ("Page 3\nCompare Models",
     ["Full 9-metric styled table\n  (teal=best, red=worst)",
      "5 accuracy bar charts",
      "4 beyond-accuracy charts:\n  Coverage / Novelty / Pop / ILD",
      "Key Takeaways summary",
      "Download results.csv button"]),
]
col_w = Inches(3.9)
sx    = Inches(0.4)
for i, (title, items) in enumerate(pages):
    bx = sx + i * (col_w + Inches(0.22))
    box(s, bx, Inches(1.45), col_w, Inches(5.6), fill_color=NAVY)
    box(s, bx, Inches(1.45), col_w, Inches(0.55), fill_color=RED)
    txt(s, title, bx + Inches(0.12), Inches(1.47),
        col_w - Inches(0.24), Inches(0.52), size=13, bold=True, color=WHITE)
    bullets(s, items, bx + Inches(0.1), Inches(2.1),
            col_w - Inches(0.2), Inches(4.7), size=11, marker="  > ")

box(s, Inches(0), Inches(6.88), W, Inches(0.62), fill_color=RGBColor(0x16, 0x21, 0x3e))
txt(s, "streamlit run app.py   ->   http://localhost:8501",
    Inches(0), Inches(6.92), W, Inches(0.45),
    size=13, color=TEAL, align=PP_ALIGN.CENTER, italic=True)


# ── 11. KEY TAKEAWAYS ────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x10, 0x15, 0x28))
slide_title(s, "Key Takeaways", "What the numbers tell us")

takeaways = [
    ("MatrixFact wins on accuracy",
     "P@10=0.090, NDCG=0.112, HR=0.498 -- leads at K=5, 10 and 20.\n50 SVD factors compress 98.3%-sparse data into dense taste signals."),
    ("MostPopular is hard to beat",
     "P@10=0.064 with zero personalisation. Coverage of 1% is the cost:\nit treats 9,742 movies as if only ~97 exist."),
    ("Accuracy vs. novelty is a real trade-off",
     "ContentBased: novelty=14.64, coverage=32% -- but P@10=0.008.\nMatrixFact: best accuracy -- but AvgPop=126, blockbuster-heavy."),
    ("ILD reveals a new dimension",
     "ContentBased ILD=0.20 -- genre-monotone despite personalisation.\nMostPopular ILD=0.95 -- popular films naturally span many genres."),
    ("ItemItemCF is constrained by sparsity",
     "Only 35.8% of items have >= 5 co-raters. The model is not broken.\nPaired with a fallback it surfaces niche items better than any other method."),
    ("Temporal split changes the numbers",
     "Using each user's most recent 20% as test (not random) eliminates\ntemporal leakage and gives a realistic picture of deployment performance."),
]
box_h = Inches(0.83)
gap   = Inches(0.06)
sy    = Inches(1.38)
for i, (title, body) in enumerate(takeaways):
    row = i // 2
    col = i  % 2
    bx  = Inches(0.4) + col * Inches(6.4)
    by  = sy + row * (box_h + gap)
    box(s, bx, by, Inches(6.1), box_h, fill_color=NAVY)
    box(s, bx, by, Inches(0.28), box_h, fill_color=RED)
    txt(s, title, bx + Inches(0.36), by + Inches(0.04),
        Inches(5.6), Inches(0.3), size=12, bold=True, color=WHITE)
    txt(s, body,  bx + Inches(0.36), by + Inches(0.34),
        Inches(5.6), Inches(0.46), size=10, color=RGBColor(0xcc, 0xcc, 0xff))


# ── 12. CONCLUSIONS ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x0d, 0x12, 0x26))
box(s, Inches(0), Inches(0), W, Inches(1.2), fill_color=NAVY)
txt(s, "Conclusions & Next Steps",
    Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
    size=30, bold=True, color=WHITE)
box(s, Inches(0.5), Inches(1.18), Inches(12.33), Pt(3), fill_color=RED)

use_cases = [
    ("Accuracy",   "Matrix Factorization", "P@10=0.090, best at all K values"),
    ("Discovery",  "Content-Based",        "32% coverage, novelty=14.64"),
    ("Cold-start", "MostPopular fallback",  "No training data needed"),
    ("Variety",    "MatrixFact + MMR",      "ILD=0.926, diverse genre spread"),
]
txt(s, "Use Case Recommendations",
    Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
    size=14, bold=True, color=TEAL)
for i, (uc, method, reason) in enumerate(use_cases):
    ry = Inches(1.78) + i * Inches(0.72)
    box(s, Inches(0.5), ry, Inches(5.5), Inches(0.65), fill_color=NAVY)
    txt(s, uc,     Inches(0.62), ry + Inches(0.05), Inches(1.3), Inches(0.3),
        size=11, bold=True, color=RED)
    txt(s, method, Inches(1.95), ry + Inches(0.05), Inches(2.2), Inches(0.3),
        size=11, bold=True, color=WHITE)
    txt(s, reason, Inches(1.95), ry + Inches(0.32), Inches(3.9), Inches(0.28),
        size=9, color=RGBColor(0xaa, 0xaa, 0xcc))

next_steps = [
    "Hybrid model: blend MatrixFact + ContentBased for accuracy + novelty",
    "Proper MF with bias terms: SGD-MF or implicit ALS (surprise / implicit library)",
    "MMR re-ranking for ContentBased: raise ILD from 0.20 without losing relevance",
    "Implicit feedback: treat views-without-ratings as weak positive signals",
]
txt(s, "Future Work",
    Inches(6.5), Inches(1.3), Inches(6.3), Inches(0.4),
    size=14, bold=True, color=TEAL)
bullets(s, next_steps, Inches(6.5), Inches(1.78), Inches(6.3), Inches(3.0),
        size=12, marker="  ->  ")

box(s, Inches(0), Inches(6.5), W, Inches(1.0), fill_color=RED)
txt(s, "CineMatch  |  ESADE Recommender Systems  |  2025/26  |  MovieLens Small  |  7 algorithms  |  9 metrics",
    Inches(0), Inches(6.62), W, Inches(0.5),
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── SAVE ─────────────────────────────────────────────────────────────────────
prs.save("slides.pptx")
print(f"Saved slides.pptx  ({len(prs.slides)} slides)")
